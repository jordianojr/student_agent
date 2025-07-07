import os
import logging
from datetime import datetime
from typing import List
from openai import OpenAI
from fastapi import UploadFile
from pymongo import MongoClient
from pptx import Presentation
import pymupdf as fitz
import hashlib
import asyncio
from dotenv import load_dotenv
from agents.semantic_chunker import SemanticChunker

_ = load_dotenv()

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# OpenAI setup
openai_client = OpenAI(api_key=os.environ["OPENAI_API_KEY"])
MAX_TOKENS = 120_000

class FileProcessor:
    def __init__(self):
        self.supported_formats = {"pdf"}
        self.db = None
        self.embedding_model = "text-embedding-3-small"
        self.embedding_dimensions = 1536
        self.semantic_chunker = SemanticChunker()    
        
    def init_db(self, mongodb_url: str = "mongodb://localhost:27017"):
        """Initialize database connection"""
        try:
            client = MongoClient(mongodb_url)
            self.db = client.agents_db
            client.admin.command('ping')
            logger.info("FileProcessor connected to MongoDB!")
        except Exception as e:
            logger.error(f"MongoDB connection error: {e}")
            raise
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 100) -> List[str]:
        """
        Split text into chunks for embedding
        
        Args:
            text: Input text to chunk
            chunk_size: Maximum characters per chunk
            overlap: Number of characters to overlap between chunks
        """
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to end at a sentence boundary
            if end < len(text):
                # Look for sentence endings within the last 100 characters
                last_period = text.rfind('.', end - 100, end)
                last_newline = text.rfind('\n', end - 100, end)
                boundary = max(last_period, last_newline)
                
                if boundary > start:
                    end = boundary + 1
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - overlap
            
        return chunks
    
    def get_embeddings(self, texts: List[str]) -> List[List[float]]:
        """
        Get embeddings for a list of texts using OpenAI API
        
        Args:
            texts: List of text strings to embed
            
        Returns:
            List of embedding vectors
        """
        try:
            response = openai_client.embeddings.create(
                model=self.embedding_model,
                input=texts
            )
            return [data.embedding for data in response.data]
        except Exception as e:
            logger.error(f"Error getting embeddings from OpenAI: {e}")
            raise
    
    def generate_chunk_id(self, filename: str, chunk_index: int, chunk_text: str) -> str:
        """Generate unique ID for a chunk"""
        content = f"{filename}_{chunk_index}_{chunk_text[:100]}"
        return hashlib.md5(content.encode()).hexdigest()
    
    def tokenize_text(self, text: str) -> dict:
        """Tokenize text using OpenAI's tokenizer (placeholder implementation)"""
        # This is a simplified tokenization - you might want to use tiktoken for accuracy
        words = text.split()
        token_count = len(words) * 1.3  # Rough approximation
        return {
            "tokens": words,
            "token_count": int(token_count)
        }
    
    def extract_notes_pptx(self, filepath):
        """Extract notes from a PowerPoint file (placeholder implementation)"""
        # This function should be implemented to extract notes from PPTX files
        try:
            prs = Presentation(filepath)
            notes = []
            for slide in prs.slides:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text
                    notes.append(notes_text)
            return "\n".join(notes) if notes else None
        except Exception as e:
            logger.error(f"Error extracting notes from PPTX: {e}")
        return None
    
    async def store_embeddings(self, filename: str, chunks: List[str], embeddings: List[List[float]], 
                             metadata: dict = None) -> int:
        """
        Store text chunks and their embeddings in MongoDB
        
        Args:
            filename: Name of the source file
            chunks: List of text chunks
            embeddings: List of embedding vectors
            metadata: Additional metadata for the document
            
        Returns:
            Number of documents inserted
        """
        documents_to_insert = []
        
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            doc = {
                "_id": self.generate_chunk_id(filename, i, chunk),
                "filename": filename,
                "chunk_index": i,
                "text": chunk,
                "embedding": embedding,
                "embedding_model": self.embedding_model,
                "created_at": datetime.utcnow(),
                "metadata": metadata or {}
            }
            documents_to_insert.append(doc)
        
        try:
            # Insert documents into the 'documents' collection
            result = self.db.semantic_documents.insert_many(documents_to_insert, ordered=False)
            return len(result.inserted_ids)
        except Exception as e:
            logger.error(f"Error inserting embeddings: {e}")
            # Try inserting one by one to handle duplicates
            inserted_count = 0
            for doc in documents_to_insert:
                try:
                    await self.db.semantic_documents.insert_one(doc)
                    inserted_count += 1
                except Exception:
                    logger.warning(f"Skipped duplicate chunk {doc['chunk_index']} for {filename}")
            return inserted_count
    
    async def process_files(self, files: List[UploadFile]) -> None:
        """Process files and update the database with embeddings"""
        if self.db is None:
            await self.init_db()
            
        processed_files = []
        self.total_tokens = 0
        
        try:
            for file in files:
                # Save file temporarily to process
                temp_path = f"/tmp/{file.filename}"
                
                try:
                    contents = await file.read()
                    with open(temp_path, "wb") as temp_file:
                        temp_file.write(contents)
                    
                    # Extract text
                    extracted_text = self.extract_pdf_text(temp_path)
                    extracted_notes = self.extract_notes_pptx(temp_path)
                    if extracted_notes:
                        extracted_text = extracted_notes if extracted_text is None else f"{extracted_text}\n{extracted_notes}"
                        logger.info(f"Content extracted from {file.filename}: {extracted_text}")
                    if extracted_text is None:
                        logger.error(f"Failed to extract text from {file.filename}")
                        continue
                    
                    # Tokenize the extracted text
                    tokenization_info = self.tokenize_text(extracted_text)
                    new_tokens = tokenization_info["token_count"]
                    
                    # Check token limit
                    if self.total_tokens + new_tokens > MAX_TOKENS:
                        logger.warning(f"Skipping {file.filename} as it would exceed token limit")
                        continue
                    
                    self.total_tokens += new_tokens
                    
                    # Chunk text for embeddings
                    # chunks = self.chunk_text(extracted_text, chunk_size=1000, overlap=100)
                    chunks = self.semantic_chunker.split_text(extracted_text)
                    logger.info(f"Created {len(chunks)} chunks for {file.filename}")
                    
                    # Get embeddings in batches
                    all_embeddings = []
                    batch_size = 100  # OpenAI API limit
                    
                    for i in range(0, len(chunks), batch_size):
                        batch_chunks = chunks[i:i + batch_size]
                        logger.info(f"Processing embedding batch {i//batch_size + 1}/{(len(chunks) + batch_size - 1)//batch_size} for {file.filename}")
                        
                        batch_embeddings = self.get_embeddings(batch_chunks)
                        all_embeddings.extend(batch_embeddings)
                        
                        # Small delay to avoid rate limiting
                        await asyncio.sleep(0.1)
                    
                    # Store embeddings in MongoDB
                    metadata = {
                        "content_type": file.content_type,
                        "total_chunks": len(chunks),
                        "original_token_count": new_tokens
                    }
                    
                    inserted_count = await self.store_embeddings(
                        file.filename, 
                        chunks, 
                        all_embeddings, 
                        metadata
                    )
                    
                    logger.info(f"Stored {inserted_count} embedding documents for {file.filename}")
                    
                    # Prepare file document for original collection (if you still need it)
                    file_doc = {
                        "filename": file.filename,
                        "content_type": file.content_type,
                        "content": {
                            "content": extracted_text,
                            "tokens": tokenization_info["tokens"],
                            "token_count": tokenization_info["token_count"]
                        },
                        "embedding_info": {
                            "total_chunks": len(chunks),
                            "embedding_model": self.embedding_model,
                            "chunks_stored": inserted_count
                        },
                        "processed_at": datetime.utcnow()
                    }
                    processed_files.append(file_doc)
                    
                finally:
                    # Clean up temp file
                    if os.path.exists(temp_path):
                        os.remove(temp_path)
                        
        except Exception as e:
            logger.error(f"Error processing files: {e}")
            raise
        
        # Store original file documents if needed
        if processed_files:
            try:
                self.db.files.insert_many(processed_files)
                logger.info(f"Stored {len(processed_files)} file documents")
            except Exception as e:
                logger.error(f"Error storing file documents: {e}")
    
    def extract_pdf_text(self, file_path):
        """Extract text from PDF file"""
        try:
            doc = fitz.open(file_path)
            text = ""
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            logger.error(f"Error extracting PDF: {e}")
            return None
    
    def clean_text(self, text):
        """
        Cleans the given text by normalizing unicode characters to ASCII, replacing newlines and tabs with a space,
        replacing multiple spaces with a single space, and trimming leading and trailing spaces.

        Args:
            text (str): The text to be cleaned.

        Returns:
            str: The cleaned text.
        """
        # Replace newlines and tabs with a space
        text = text.replace("\n", " ").replace("\t", " ")

        # Replace multiple spaces with a single space
        text = " ".join(text.split())

        # Trim leading and trailing spaces
        text = text.strip()

        return text

    def vector_search(self, query: str, filename_filter: str = None) -> List[dict]:
        """
        Perform vector search using MongoDB Atlas Vector Search
        
        Args:
            query: Search query text
            limit: Number of results to return
            filename_filter: Optional filter by filename
            
        Returns:
            List of matching documents with scores
        """
        mongodb_url = os.environ.get("MONGODB_URL")
        if self.db is None:
            self.init_db(mongodb_url=mongodb_url)
        
        # Get query embedding
        query_embedding = self.get_embeddings([query])[0]
        
        # Build aggregation pipeline
        pipeline = [
            {
                "$vectorSearch": {
                    "index": "vector_index_2",  # Your vector index name
                    "path": "embedding",      # Field containing the vectors
                    "queryVector": query_embedding,  # Your query vector (list of floats)
                    "numCandidates": 100,     # Number of candidates to consider
                    "limit": 10               # Number of results to return
                }
            },
            {
                "$project": {
                    "filename": 1,        # Include the page content
                    "text": 1,            # Include text content
                    "chunk_index": 1,    # Include chunk index
                    "embedding": 1,           # Include embedding if needed
                    "score": {                # Include similarity score
                        "$meta": "vectorSearchScore"
                    }
                }
            }
        ]
        
        # Add filter if specified
        if filename_filter:
            pipeline[0]["$vectorSearch"]["filter"] = {"filename": {"$in": filename_filter}}
        
        try:
            result = list(self.db.semantic_documents.aggregate(pipeline))
            # Step 3: Return the top result's chunk
            if result:
                print(f"Found {len(result)} matching documents.")
                print(f"Top result: {result[0]["chunk_index"]}")
                result = result[0].get("text", "No content found in top result.")
                cleaned_text = self.clean_text(result)
                return cleaned_text
            else:
                return "No matching documents found."
        except Exception as e:
            logger.error(f"Vector search error: {e}")
            return []
    
    async def get_document_stats(self, filename: str = None) -> dict:
        """Get statistics about stored documents"""
        match_filter = {}
        if filename:
            match_filter["filename"] = filename
        
        pipeline = [
            {"$match": match_filter},
            {"$group": {
                "_id": "$filename",
                "chunk_count": {"$sum": 1},
                "avg_text_length": {"$avg": {"$strLenCP": "$text"}},
                "created_at": {"$first": "$created_at"}
            }}
        ]
        
        cursor = self.db.semantic_documents.aggregate(pipeline)
        stats = await cursor.to_list(length=None)
        
        return {
            "total_documents": len(stats),
            "documents": stats,
            "total_chunks": sum(doc["chunk_count"] for doc in stats)
        }
    
    async def delete_document_embeddings(self, filename: str) -> int:
        """Delete all chunks for a specific file"""
        result = await self.db.semantic_documents.delete_many({"filename": filename})
        return result.deleted_count

# Create a singleton instance
file_processor = FileProcessor()